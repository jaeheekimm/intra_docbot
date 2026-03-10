import os
import re
from typing import Dict, List, Any, Optional, Tuple

from pptx import Presentation
from langchain_core.documents import Document

from src.utils.files import safe_filename
from src.utils.hashing import sha1_short


def _normalize_pptx_text(t: str) -> str:
    """PPTX 텍스트의 줄바꿈/공백을 가볍게 정리"""
    s = (t or "").replace("\r\n", "\n").replace("\r", "\n")
    s = re.sub(r"[ \t]{2,}", " ", s).strip()
    return s


def _extract_sorted_text_blocks(slide) -> List[Tuple[int, int, str]]:
    """슬라이드 내 텍스트 프레임을 (top, left) 좌표 기준으로 정렬해서 반환.

    반환: [(top, left, text), ...]
    """
    blocks: List[Tuple[int, int, str]] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            raw = shape.text_frame.text
            text = _normalize_pptx_text(raw)
            if not text:
                continue
            top = int(getattr(shape, "top", 0) or 0)
            left = int(getattr(shape, "left", 0) or 0)
            blocks.append((top, left, text))

    blocks.sort(key=lambda x: (x[0], x[1]))  # 위→아래, 왼→오
    return blocks


def _pick_title_from_blocks(
    blocks: List[Tuple[int, int, str]],
) -> Tuple[Optional[str], List[str]]:
    """정렬된 블록 중 상단 텍스트를 보수적으로 title로 추정.

    - 상위 5개 블록 중 80자 이하 텍스트를 title로 선택
    - 못 고르면 None 반환
    반환: (title or None, body_texts)
    """
    if not blocks:
        return None, []

    texts = [b[2] for b in blocks]

    title: Optional[str] = None
    title_idx: Optional[int] = None
    for i, t in enumerate(texts[:5]):
        one_line = t.replace("\n", " ").strip()
        if 1 <= len(one_line) <= 80:
            title = one_line
            title_idx = i
            break

    if title_idx is None:
        return None, texts

    body = texts[:title_idx] + texts[title_idx + 1 :]
    return title, body


def extract_pptx_images(pptx_path: str, out_dir: str) -> Dict[int, List[str]]:
    """PPTX에서 이미지 추출 후 out_dir에 저장.

    반환: {슬라이드번호(1-based): [이미지경로, ...]}
    - 1번 슬라이드 이미지 전부 스킵 (표지 로고 등)
    - 2번 슬라이드부터는 오른쪽 위 코너 이미지만 스킵 (헤더 로고)
      → 코너 기준: 가로 82%, 세로 18% 안쪽 영역
    """
    mapping: Dict[int, List[str]] = {}

    prs = Presentation(pptx_path)
    base = safe_filename(os.path.basename(pptx_path))

    sw, sh = prs.slide_width, prs.slide_height
    RU_X = int(sw * 0.82)  # 오른쪽 위 코너 기준 x (비율 조정 시 여기만)
    RU_Y = int(sh * 0.18)  # 오른쪽 위 코너 기준 y

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) == 13 and getattr(
                shape, "image", None
            ):
                x, y = int(shape.left), int(shape.top)

                if slide_idx == 1:
                    continue

                if slide_idx >= 2 and (x >= RU_X) and (y <= RU_Y):
                    continue

                img = shape.image
                blob = img.blob
                ext = img.ext

                hid = sha1_short(blob)
                fname = f"{base}__slide{slide_idx:03d}__{hid}.{ext}"
                save_path = os.path.join(out_dir, fname)

                if not os.path.exists(save_path):
                    with open(save_path, "wb") as f:
                        f.write(blob)

                mapping.setdefault(slide_idx, []).append(save_path)

    return mapping


def load_pptx_docs(pptx_path: str, ppt_img_map: Dict[int, List[str]]) -> List[Document]:
    """PPTX를 슬라이드 단위 Document 리스트로 변환.

    - 텍스트는 (top, left) 좌표 기준으로 정렬
    - title은 metadata에 별도 저장 (본문에서 제거하지 않음)
    - 텍스트도 이미지도 없는 슬라이드는 스킵
    """
    prs = Presentation(pptx_path)
    file_name = os.path.basename(pptx_path)

    docs: List[Document] = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        blocks = _extract_sorted_text_blocks(slide)
        title, body_texts = _pick_title_from_blocks(blocks)

        content = "\n".join([t for t in body_texts if t]).strip()

        if not title and not content and not ppt_img_map.get(slide_idx):
            continue

        md = {
            "source": pptx_path,
            "file_name": file_name,
            "doc_type": "pptx",
            "page": None,
            "slide": slide_idx,
            "sheet": None,
            "row": None,
            "title": title,
            "image_paths": ppt_img_map.get(slide_idx, []),
        }

        docs.append(Document(page_content=content, metadata=md))

    return docs


def pptx_image_manifest(
    pptx_path: str, ppt_img_map: Dict[int, List[str]]
) -> List[Dict[str, Any]]:
    """슬라이드별 이미지 경로 정보를 image_manifest 포맷의 딕셔너리 리스트로 반환"""
    items: List[Dict[str, Any]] = []
    for slide, paths in ppt_img_map.items():
        for ip in paths:
            items.append(
                {
                    "source_file": pptx_path,
                    "source_type": "pptx",
                    "page": None,
                    "slide": slide,
                    "sheet": None,
                    "row": None,
                    "image_path": ip,
                }
            )
    return items
